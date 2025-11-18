import os
import shutil
import sys
import json
import threading
import webbrowser
import customtkinter as ctk
import sqlite3
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.backends import default_backend
from cryptography.hazmat.primitives import hashes
from base64 import urlsafe_b64encode, urlsafe_b64decode
from tkinter import messagebox, StringVar, filedialog, Listbox, END
from docx import Document
from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
from odf.text import P
from odf.table import Table, TableRow, TableCell
from odf.draw import Page, Frame, TextBox, Image
from odf.style import MasterPage
from PIL import Image, ImageTk
import urllib
import ctypes
import queue
import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from ctypes import *
from ctypes.wintypes import *
if sys.platform.startswith("win"):
    import winreg
    user32 = windll.user32
    shell32 = windll.shell32
    kernel32 = windll.kernel32
    WM_USER = 0x0400
    WM_COMMAND = 0x0111
    WM_LBUTTONDBLCLK = 0x0203
    WM_RBUTTONDOWN = 0x0204
    NIM_ADD = 0
    NIM_MODIFY = 1
    NIM_DELETE = 2
    NIF_MESSAGE = 0x00000001
    NIF_ICON = 0x00000002
    NIF_TIP = 0x00000004
    NIF_INFO = 0x00000010
    NIIF_INFO = 0x00000001
    IDI_APPLICATION = 32512
    MF_STRING = 0x00000000
    MF_POPUP = 0x00000010
    IMAGE_ICON = 1
    LR_LOADFROMFILE = 0x00000010
    LR_DEFAULTSIZE = 0x00000040
    class POINT(Structure):
        _fields_ = [("x", LONG), ("y", LONG)]
    class NOTIFYICONDATA(Structure):
        _fields_ = [
            ("cbSize", DWORD),
            ("hWnd", HWND),
            ("uID", UINT),
            ("uFlags", UINT),
            ("uCallbackMessage", UINT),
            ("hIcon", HICON),
            ("szTip", c_char * 64),
            ("dwState", DWORD),
            ("dwStateMask", DWORD),
            ("szInfo", c_char * 256),
            ("uVersion", UINT),
            ("szInfoTitle", c_char * 64),
            ("dwInfoFlags", DWORD),
        ]
    HWND = c_void_p
    UINT = c_uint
    WPARAM = c_ulonglong
    LPARAM = c_longlong
    WNDPROC = WINFUNCTYPE(c_longlong, HWND, UINT, WPARAM, LPARAM)
    user32.CallWindowProcA.argtypes = [c_void_p, HWND, UINT, WPARAM, LPARAM]
    user32.CallWindowProcA.restype = c_longlong
    # Added argtypes/restype for single-instance handling
    user32.FindWindowW.argtypes = [c_wchar_p, c_wchar_p]
    user32.FindWindowW.restype = HWND
    user32.PostMessageW.argtypes = [HWND, UINT, WPARAM, LPARAM]
    user32.PostMessageW.restype = BOOL
    # Definitions for missing types moved here before usage
    class SECURITY_ATTRIBUTES(Structure):
        _fields_ = [
            ("nLength", DWORD),
            ("lpSecurityDescriptor", LPVOID),
            ("bInheritHandle", BOOL),
        ]
    LPSECURITY_ATTRIBUTES = POINTER(SECURITY_ATTRIBUTES)
    UINT_PTR = c_ulonglong
    HMENU = c_void_p
    LPVOID = c_void_p
    HINSTANCE = c_void_p
    HICON = c_void_p
    HANDLE = c_void_p
    BOOL = c_int
    kernel32.CreateMutexW.argtypes = [LPSECURITY_ATTRIBUTES, BOOL, c_wchar_p]
    kernel32.CreateMutexW.restype = HANDLE
    kernel32.GetLastError.argtypes = []
    kernel32.GetLastError.restype = DWORD
    user32.GetWindowLongPtrA.argtypes = [HWND, c_int]
    user32.GetWindowLongPtrA.restype = c_longlong
    user32.SetWindowLongPtrA.argtypes = [HWND, c_int, c_longlong]
    user32.SetWindowLongPtrA.restype = c_longlong
    user32.RegisterWindowMessageA.argtypes = [c_char_p]
    user32.RegisterWindowMessageA.restype = UINT
    shell32.Shell_NotifyIconA.argtypes = [DWORD, POINTER(NOTIFYICONDATA)]
    shell32.Shell_NotifyIconA.restype = BOOL
    user32.LoadImageA.argtypes = [HINSTANCE, c_char_p, UINT, c_int, c_int, UINT]
    user32.LoadImageA.restype = HANDLE
    user32.CreatePopupMenu.argtypes = []
    user32.CreatePopupMenu.restype = HMENU
    user32.AppendMenuA.argtypes = [HMENU, UINT, UINT_PTR, c_char_p]
    user32.AppendMenuA.restype = BOOL
    user32.GetCursorPos.argtypes = [POINTER(POINT)]
    user32.GetCursorPos.restype = BOOL
    user32.SetForegroundWindow.argtypes = [HWND]
    user32.SetForegroundWindow.restype = BOOL
    user32.TrackPopupMenu.argtypes = [HMENU, UINT, c_int, c_int, c_int, HWND, c_void_p]
    user32.TrackPopupMenu.restype = BOOL
    user32.PostMessageA.argtypes = [HWND, UINT, WPARAM, LPARAM]
    user32.PostMessageA.restype = BOOL
    user32.DestroyMenu.argtypes = [HMENU]
    user32.DestroyMenu.restype = BOOL
# Single instance enforcement using mutex
if sys.platform.startswith("win"):
    ERROR_ALREADY_EXISTS = 183
    mutex_name = "Global\\BlackHole_SingleInstance_Mutex"
    mutex = kernel32.CreateMutexW(None, True, mutex_name)
    err = kernel32.GetLastError()
    if err == ERROR_ALREADY_EXISTS:
        hwnd = user32.FindWindowW(None, "Black Hole Password Manager")
        if hwnd:
            user32.PostMessageW(hwnd, WM_USER + 1, 0, WM_LBUTTONDBLCLK)
        sys.exit(0)
# Set working directory to the script/exe directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(sys.argv[0]))
os.chdir(SCRIPT_DIR)
# --- App Icon ---
APP_ICON_PATH = os.path.join(SCRIPT_DIR, "Icons", "BlackHole_Icon.ico")
BLACK_HOLE_LOGO = os.path.join(SCRIPT_DIR, "Icons", "BlackHole_Transparent_Light.png")
BLACK_HOLE_WIDE_LOGO = os.path.join(SCRIPT_DIR, "Icons", "BlackHole_Transparent_Wide.png")
NOVA_FOUNDRY_LOGO = os.path.join(SCRIPT_DIR, "Icons", "Nova_foundry_wide_transparent.png")
FONT_REGULAR = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-Regular.ttf")
FONT_MEDIUM = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-Medium.ttf")
FONT_BOLD = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-Black.ttf")
FONT_LIGHT = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-Light.ttf")
FONT_ITALIC = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-Italic.ttf")
FONT_SEMIBOLD = os.path.join(SCRIPT_DIR, "Fonts", "Nunito-SemiBold.ttf")
LICENSE_TEXT = os.path.join(SCRIPT_DIR, "LICENSE.txt")
VERSION = "1.5.3"
# Load all the font files for Tkinter (on Windows)
if sys.platform.startswith("win"):
    fonts = [FONT_REGULAR, FONT_MEDIUM, FONT_BOLD, FONT_LIGHT, FONT_ITALIC, FONT_SEMIBOLD]
    for font_path in fonts:
        ctypes.windll.gdi32.AddFontResourceA(font_path.encode('utf-8'))
    # Broadcast font change
    HWND_BROADCAST = 0xFFFF
    WM_FONTCHANGE = 0x001D
    ctypes.windll.user32.SendMessageA(HWND_BROADCAST, WM_FONTCHANGE, 0, 0)
# --- Paths ---
local_appdata = os.getenv("LOCALAPPDATA") or os.getenv("APPDATA")
nova_folder = os.path.join(local_appdata, "NovaFoundry")
os.makedirs(nova_folder, exist_ok=True)
stored_icons_path = os.path.join(nova_folder, "StoredIcons")
os.makedirs(stored_icons_path, exist_ok=True)
settings_path = os.path.join(nova_folder, "settings.json")
order_path = os.path.join(nova_folder, "order.json")
pinned_path = os.path.join(nova_folder, "pinned.json")
# --- Theme (Deep Space Glow) colours ---
BG = "#05050a"
CARD = "#0b0b0f"
CARD_HOVER = "#111327"
ACCENT = "#47a3ff"
ACCENT_DIM = "#2b6f9f"
TEXT = "#e6eef8"
# --- Tooltip Class (Fixed version with delay to prevent sticking) ---
class Tooltip:
    def __init__(self, widget, text, wait_time=500):
        self.widget = widget
        self.text = text
        self.wait_time = wait_time
        self.tooltip = None
        self.id = None
        self.widget.bind("<Enter>", self.schedule_show)
        self.widget.bind("<Leave>", self.hide_tooltip)
    def schedule_show(self, event):
        if self.id is not None:
            self.widget.after_cancel(self.id)
        self.id = self.widget.after(self.wait_time, lambda: self.show_tooltip(event))
    def show_tooltip(self, event):
        self.id = None
        if self.tooltip:
            return
        x = event.x_root + 20
        y = event.y_root + 20
        self.tooltip = ctk.CTkToplevel(self.widget)
        self.tooltip.wm_overrideredirect(True)
        self.tooltip.wm_geometry(f"+{x}+{y}")
        label = ctk.CTkLabel(self.tooltip, text=self.text, corner_radius=0, fg_color=CARD, text_color=TEXT, padx=10, pady=5)
        label.pack()
    def hide_tooltip(self, event):
        if self.id is not None:
            self.widget.after_cancel(self.id)
            self.id = None
        if self.tooltip:
            self.tooltip.destroy()
            self.tooltip = None
# --- Helper: Derive Key from master password ---
def derive_key(password: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=390000,
        backend=default_backend()
    )
    return urlsafe_b64encode(kdf.derive(password.encode()))
# --- Helper: Center Popup ---
def center_popup(popup):
    popup.update_idletasks()
    width = popup.winfo_reqwidth()
    height = popup.winfo_reqheight()
    x = (popup.winfo_screenwidth() // 2) - (width // 2)
    y = (popup.winfo_screenheight() // 2) - (height // 2)
    popup.geometry(f"{width}x{height}+{x}+{y}")
# --- Password Manager App ---
class PasswordManager(ctk.CTk):
    def __init__(self):
        super().__init__()
        ctk.set_appearance_mode("dark")
        try:
            ctk.set_default_color_theme("dark-blue")
        except Exception:
            pass
        # Hide main window until auth succeeds
        self.title("Black Hole Password Manager")
        self.geometry("900x560")
        if os.path.exists(APP_ICON_PATH):
            try:
                self.iconbitmap(APP_ICON_PATH)
            except Exception:
                pass
        self.attributes('-alpha', 0.0)
        # encryption / state
        self.master_password = None
        self.fernet = None
        self.salt = None
        self.db_path = None
        self.conn = None
        self.c = None
        self.authenticated = False
        self.ui_built = False
        # load settings
        self.settings = {"master_password_set": False}
        if os.path.exists(settings_path):
            try:
                with open(settings_path, "r") as f:
                    self.settings = json.load(f)
                if "salt" in self.settings:
                    self.salt = urlsafe_b64decode(self.settings["salt"])
                if "db_path" in self.settings:
                    self.db_path = self.settings["db_path"]
            except Exception:
                self.settings = {"master_password_set": False}
        # load order settings
        self.order_mode = "default"
        self.custom_order = []
        if os.path.exists(order_path):
            try:
                with open(order_path, "r") as f:
                    data = json.load(f)
                    self.order_mode = data.get("mode", "default")
                    self.custom_order = data.get("custom_order", [])
            except Exception:
                pass
        # load pinned
        self.pinned = []
        if os.path.exists(pinned_path):
            try:
                with open(pinned_path, "r") as f:
                    self.pinned = json.load(f).get("pinned", [])
            except Exception:
                pass
        self.configure(fg_color=BG)
        # Initialize DB early (titles are plaintext)
        if self.db_path:
            self._init_db()
        # Check if setup is needed
        if not self.settings.get("master_password_set", False) or not self.db_path:
            success = self._show_setup_modal()
            if success:
                self.authenticated = True
        else:
            if not ('--minimize' in sys.argv and self.settings.get("minimize_to_tray", False)):
                success = self._show_master_password_modal()
                if success:
                    self.authenticated = True
            else:
                success = True # Delay auth
                self.authenticated = False
        if not success:
            self.destroy()
            sys.exit()
        if self.authenticated:
            self._build_ui()
            self.attributes('-alpha', 1.0)
            self.after(0, lambda: self.state('zoomed'))
            self.load_cards()
            self.ui_built = True
        # Tray setup
        if sys.platform.startswith("win"):
            self.hwnd = None
            self.callback_message = None
            self.new_wndproc_ptr = None
            self.original_wndproc = None
            self.msg_taskbar_created = None
            self.tray_added = False
            self.protocol("WM_DELETE_WINDOW", self.on_close)
            self.message_queue = queue.Queue()
            self.after(100, self.process_message_queue)
            self.add_tray() # Always add tray to persist
        # Minimize on start if applicable
        if '--minimize' in sys.argv and self.settings.get("minimize_to_tray", False):
            self.withdraw()
    def process_message_queue(self):
        while not self.message_queue.empty():
            msg_type, data = self.message_queue.get()
            if msg_type == 'tray_msg':
                self.handle_tray_msg(data)
            elif msg_type == 'menu_cmd':
                self.handle_menu_cmd(data)
            elif msg_type == 'taskbar_created':
                self.remove_tray()
                self.add_tray()
        self.after(100, self.process_message_queue)
    def on_close(self):
        if self.settings.get("minimize_to_tray", False):
            self.withdraw()
        else:
            self.exit_app()
    def exit_app(self):
        self.remove_tray()
        self.destroy()
    # --- Tray functions ---
    def add_tray(self):
        if not sys.platform.startswith("win") or self.tray_added:
            return
        self.hwnd = self.winfo_id()
        self.msg_taskbar_created = user32.RegisterWindowMessageA(b"TaskbarCreated")
        self.callback_message = WM_USER + 1
        def new_wndproc(hwnd, msg, wparam, lparam):
            if msg == self.callback_message:
                self.message_queue.put(('tray_msg', lparam))
                return 0
            elif msg == WM_COMMAND:
                cmd = wparam & 0xFFFF
                self.message_queue.put(('menu_cmd', cmd))
                return 0
            elif msg == self.msg_taskbar_created:
                self.message_queue.put(('taskbar_created', None))
                return 0
            return user32.CallWindowProcA(c_void_p(self.original_wndproc), hwnd, msg, wparam, lparam)
        self.new_wndproc_ptr = WNDPROC(new_wndproc)
        self.original_wndproc = user32.GetWindowLongPtrA(self.hwnd, -4)
        user32.SetWindowLongPtrA(self.hwnd, -4, c_longlong(cast(self.new_wndproc_ptr, c_void_p).value))
        nid = NOTIFYICONDATA()
        nid.cbSize = sizeof(NOTIFYICONDATA)
        nid.hWnd = self.hwnd
        nid.uID = 99
        nid.uFlags = NIF_ICON | NIF_MESSAGE | NIF_TIP
        nid.uCallbackMessage = self.callback_message
        nid.szTip = b"Black Hole Password Manager\0"
        nid.hIcon = user32.LoadImageA(0, APP_ICON_PATH.encode('utf-8'), IMAGE_ICON, 0, 0, LR_LOADFROMFILE | LR_DEFAULTSIZE)
        shell32.Shell_NotifyIconA(NIM_ADD, byref(nid))
        self.tray_added = True
    def remove_tray(self):
        if not sys.platform.startswith("win") or not self.tray_added:
            return
        nid = NOTIFYICONDATA()
        nid.cbSize = sizeof(NOTIFYICONDATA)
        nid.hWnd = self.hwnd
        nid.uID = 99
        shell32.Shell_NotifyIconA(NIM_DELETE, byref(nid))
        user32.SetWindowLongPtrA(self.hwnd, -4, self.original_wndproc)
        self.tray_added = False
    def handle_tray_msg(self, lparam):
        if lparam == WM_LBUTTONDBLCLK:
            self.restore_from_tray()
        elif lparam == WM_RBUTTONDOWN:
            self.show_tray_menu()
    def handle_menu_cmd(self, cmd):
        if cmd == 1001:
            self.restore_from_tray()
        elif cmd == 1002:
            self.exit_app()
        elif cmd >= 2000 and cmd < 2000 + len(self.pinned):
            id_ = self.pinned[cmd - 2000]
            self.copy_pinned(id_)
    def show_tray_menu(self):
        menu = user32.CreatePopupMenu()
        pinned_menu = user32.CreatePopupMenu()
        # Pinned items
        for idx, id_ in enumerate(self.pinned):
            row = self.c.execute("SELECT title FROM passwords WHERE id=?", (id_,)).fetchone()
            if row:
                title = row[0][:50].encode('utf-8') # Truncate if too long
                user32.AppendMenuA(pinned_menu, MF_STRING, 2000 + idx, title)
        if len(self.pinned) == 0:
            user32.AppendMenuA(pinned_menu, MF_STRING, 0, b"No pinned accounts")
        user32.AppendMenuA(menu, MF_STRING | MF_POPUP, pinned_menu, b"Pinned Accounts")
        user32.AppendMenuA(menu, MF_STRING, 1001, b"Open")
        user32.AppendMenuA(menu, MF_STRING, 1002, b"Exit")
        pt = POINT()
        user32.GetCursorPos(byref(pt))
        user32.SetForegroundWindow(self.hwnd)
        user32.TrackPopupMenu(menu, 0, pt.x, pt.y, 0, self.hwnd, None)
        user32.PostMessageA(self.hwnd, 0, 0, 0)
        user32.DestroyMenu(menu)
        user32.DestroyMenu(pinned_menu)
    def restore_from_tray(self):
        if not self.authenticated:
            success = self._show_master_password_modal()
            if not success:
                return
            self.authenticated = True
        if not self.ui_built:
            self._build_ui()
            self.load_cards()
            self.ui_built = True
        self.attributes('-alpha', 1.0)
        self.after(0, lambda: self.state('zoomed'))
        self.deiconify()
        self.lift()
    def copy_pinned(self, id_):
        if self._verify_master_password():
            row = self.c.execute("SELECT title, password FROM passwords WHERE id=?", (id_,)).fetchone()
            if row:
                title, pwd_enc = row
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
                if messagebox.askyesno("Confirm Copy", f"Are you sure you want to copy the password for '{title}' to the clipboard?"):
                    self.clipboard_clear()
                    self.clipboard_append(pwd)
                    self.show_balloon("Copied", "Password copied to clipboard!")
    def show_balloon(self, title, msg):
        nid = NOTIFYICONDATA()
        nid.cbSize = sizeof(NOTIFYICONDATA)
        nid.hWnd = self.hwnd
        nid.uID = 99
        nid.uFlags = NIF_INFO
        nid.uVersion = 3
        nid.szInfoTitle = title.encode('utf-8')[:64]
        nid.szInfo = msg.encode('utf-8')[:256]
        nid.dwInfoFlags = NIIF_INFO
        shell32.Shell_NotifyIconA(NIM_MODIFY, byref(nid))
    # --- Setup modal: New or Import ---
    def _show_setup_modal(self):
        popup = ctk.CTkToplevel(self)
        popup.title("Black Hole — Setup")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        closed_by_user = {"val": True}
        def on_close():
            closed_by_user["val"] = True
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        # Header
        ctk.CTkLabel(popup, text="Black Hole — Setup",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        ctk.CTkLabel(popup, text="Set up a new vault or import an existing one?",
                     text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0,12))
        # Buttons
        btn_frame = ctk.CTkFrame(popup, fg_color=BG, corner_radius=0)
        btn_frame.pack(pady=(12,12))
        def setup_new():
            closed_by_user["val"] = False
            popup.grab_release()
            popup.destroy()
            self._setup_new()
        def setup_import():
            closed_by_user["val"] = False
            popup.grab_release()
            popup.destroy()
            self._setup_import()
        ctk.CTkButton(btn_frame, text="New Vault", command=setup_new,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        ctk.CTkButton(btn_frame, text="Import Vault", command=setup_import,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        center_popup(popup)
        self.wait_window(popup)
        return not closed_by_user["val"]
    # --- Setup New ---
    def _setup_new(self):
        dir_path = filedialog.askdirectory(title="Select Folder for New Vault")
        if not dir_path:
            return False
        self.db_path = os.path.join(dir_path, "BlackHolePasswords.db")
        if os.path.exists(self.db_path):
            if not messagebox.askyesno("Overwrite?", "Database already exists. Overwrite?"):
                return False
        success = self._show_master_create_modal()
        if success:
            self._init_db()
            sync_key = urlsafe_b64encode(self.salt).decode()
            self._show_sync_key_display_popup(sync_key)
            try:
                with open(settings_path, "w", encoding="utf-8") as sf:
                    json.dump(self.settings, sf)
            except Exception:
                pass
        return success
    # --- Setup Import ---
    def _setup_import(self):
        file_path = filedialog.askopenfilename(title="Select Existing Vault DB", filetypes=[("Database Files", "*.db")])
        if not file_path:
            return False
        self.db_path = file_path
        # Prompt for sync key
        key_success = self._show_sync_key_modal()
        if not key_success:
            return False
        # Show master unlock modal
        unlock_success = self._show_master_unlock_modal()
        if not unlock_success:
            return False
        # Connect to DB and verify if possible
        self._init_db()
        try:
            row = self.c.execute("SELECT password FROM passwords WHERE password IS NOT NULL LIMIT 1").fetchone()
            if row and row[0]:
                self.fernet.decrypt(row[0].encode())
        except Exception:
            messagebox.showerror("Error", "Incorrect sync key or master password!")
            self.conn.close()
            self.conn = None
            self.c = None
            return False
        # Generate local verification
        self.settings["verification"] = self.fernet.encrypt(b"VERIFICATION").decode()
        self.settings["salt"] = urlsafe_b64encode(self.salt).decode()
        self.settings["db_path"] = self.db_path
        self.settings["master_password_set"] = True
        try:
            with open(settings_path, "w", encoding="utf-8") as sf:
                json.dump(self.settings, sf)
        except Exception:
            pass
        return True
    # --- Sync Key input modal ---
    def _show_sync_key_modal(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.title("Black Hole - Sync Key")
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        closed_by_user = {"val": True}
        def on_close():
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        ctk.CTkLabel(popup, text="Enter Sync Key",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        sync_var = StringVar()
        sync_entry = ctk.CTkEntry(frame, placeholder_text="Sync Key (base64)",
                                  textvariable=sync_var, width=360, justify="center")
        sync_entry.pack(padx=12, pady=(12,6))
        def submit():
            try:
                self.salt = urlsafe_b64decode(sync_var.get())
                closed_by_user["val"] = False
                popup.grab_release()
                popup.destroy()
            except Exception:
                messagebox.showerror("Error", "Invalid sync key!", parent=popup)
        ctk.CTkButton(frame, text="Submit", command=submit,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM).pack(pady=(0,10))
        sync_entry.focus_set()
        center_popup(popup)
        self.wait_window(popup)
        return not closed_by_user["val"]
    # --- Sync Key display popup ---
    def _show_sync_key_display_popup(self, sync_key):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Black Hole - Sync Key")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        def on_close():
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        ctk.CTkLabel(popup, text="Your Sync Key",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        ctk.CTkLabel(popup, text="Save this key securely to import on other devices.",
                     text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0,12))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        key_var = StringVar(value=sync_key)
        key_entry = ctk.CTkEntry(frame, textvariable=key_var, width=360, justify="center", state="readonly")
        key_entry.pack(padx=12, pady=(12,6))
        def copy_key():
            self.clipboard_clear()
            self.clipboard_append(sync_key)
            messagebox.showinfo("Copied", "Sync key copied to clipboard!", parent=popup)
        ctk.CTkButton(frame, text="Copy to Clipboard", command=copy_key,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM).pack(pady=(0,10))
        btn_frame = ctk.CTkFrame(popup, fg_color=BG, corner_radius=0)
        btn_frame.pack(pady=(12,12))
        ctk.CTkButton(btn_frame, text="OK", command=on_close,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        center_popup(popup)
        self.wait_window(popup)
    # --- Master Create modal ---
    def _show_master_create_modal(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Black Hole - Create Master Password")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        closed_by_user = {"val": True}
        def on_close():
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        ctk.CTkLabel(popup, text="Black Hole — Master Password",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        ctk.CTkLabel(popup, text="Create a master password for your vault",
                     text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0,12))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        pwd_var = StringVar()
        pwd_entry = ctk.CTkEntry(frame, placeholder_text="Master Password",
                                 show="*", textvariable=pwd_var,
                                 width=360, justify="center")
        pwd_entry.pack(padx=12, pady=(12,6))
        original_entry_color = pwd_entry.cget("fg_color")
        def toggle_pwd():
            pwd_entry.configure(show="" if pwd_entry.cget("show")=="*" else "*")
        ctk.CTkButton(frame, text="Show/Hide", width=120,
                       command=toggle_pwd, fg_color=ACCENT, text_color=BG,
                       hover_color=ACCENT_DIM).pack(pady=(0,10))
        # Optional save path for master pw
        path_var = StringVar()
        def browse_path():
            file_path = filedialog.asksaveasfilename(defaultextension=".txt",
                                                     filetypes=[("Text files","*.txt")])
            if file_path:
                path_var.set(file_path)
        ctk.CTkButton(frame, text="Select Save Path (optional)", command=browse_path,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM).pack(pady=(0,8))
        ctk.CTkLabel(frame, textvariable=path_var, text_color=TEXT, fg_color=CARD).pack(pady=(0,8))
        btn_frame = ctk.CTkFrame(popup, fg_color=BG, corner_radius=0)
        btn_frame.pack(pady=(12,12))
        def mark_wrong():
            pwd_entry.configure(fg_color="#7a2d2d")
            try:
                geo_string = popup.geometry()
                parts = geo_string.split('+')
                size_part = parts[0]
                original_x = int(parts[1])
                original_y = int(parts[2])
            except Exception:
                return
            def shake_animation(step=0):
                try:
                    offsets = [10, -10, 10, -10, 5, -5, 0]
                    if step < len(offsets):
                        dx = offsets[step]
                        popup.geometry(f"{size_part}+{original_x + dx}+{original_y}")
                        popup.after(50, shake_animation, step + 1)
                    elif step == len(offsets):
                        popup.after(1000, lambda: pwd_entry.configure(fg_color=original_entry_color))
                except Exception:
                    pass
            shake_animation(0)
        def create_master():
            pwd = pwd_var.get() or ""
            if not pwd:
                mark_wrong()
                messagebox.showerror("Error", "Master password required!", parent=popup)
                return
            try:
                self.salt = os.urandom(16)
                key = derive_key(pwd, self.salt)
                self.fernet = Fernet(key)
                verif_enc = self.fernet.encrypt(b"VERIFICATION").decode()
                self.settings["salt"] = urlsafe_b64encode(self.salt).decode()
                self.settings["verification"] = verif_enc
                self.settings["db_path"] = self.db_path
                self.settings["master_password_set"] = True
                self.master_password = pwd
            except Exception as e:
                mark_wrong()
                messagebox.showerror("Error", f"Failed to derive key: {e}", parent=popup)
                return
            save_path = path_var.get()
            if save_path:
                try:
                    with open(save_path,"w",encoding="utf-8") as f:
                        f.write(pwd)
                    messagebox.showinfo("Saved", f"Master password saved to {save_path}", parent=popup)
                except Exception as e:
                    messagebox.showwarning("Warning", f"Couldn't save: {e}", parent=popup)
            closed_by_user["val"] = False
            popup.grab_release()
            popup.destroy()
        ctk.CTkButton(btn_frame, text="Create & Continue", command=create_master,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        ctk.CTkButton(btn_frame, text="Cancel", command=on_close, fg_color="#3a3a3a", width=120).pack(side="left", padx=6)
        pwd_entry.focus_set()
        pwd_entry.bind("<Return>", lambda e: create_master())
        center_popup(popup)
        self.wait_window(popup)
        return not closed_by_user["val"] and self.fernet is not None
    # --- Master Unlock modal ---
    def _show_master_unlock_modal(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.title("Black Hole - Master Password")
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH)))
            except Exception:
                pass
        closed_by_user = {"val": True}
        def on_close():
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        ctk.CTkLabel(popup, text="Black Hole — Master Password",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        ctk.CTkLabel(popup, text="Enter your master password to unlock",
                     text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0,12))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        pwd_var = StringVar()
        pwd_entry = ctk.CTkEntry(frame, placeholder_text="Master Password",
                                 show="*", textvariable=pwd_var,
                                 width=360, justify="center")
        pwd_entry.pack(padx=12, pady=(12,6))
        original_entry_color = pwd_entry.cget("fg_color")
        def toggle_pwd():
            pwd_entry.configure(show="" if pwd_entry.cget("show")=="*" else "*")
        ctk.CTkButton(frame, text="Show/Hide", width=120,
                       command=toggle_pwd, fg_color=ACCENT, text_color=BG,
                       hover_color=ACCENT_DIM).pack(pady=(0,10))
        btn_frame = ctk.CTkFrame(popup, fg_color=BG, corner_radius=0)
        btn_frame.pack(pady=(12,12))
        pwd_entry.bind("<Return>", lambda event: unlock_master())
        def mark_wrong():
            pwd_entry.configure(fg_color="#7a2d2d")
            try:
                geo_string = popup.geometry()
                parts = geo_string.split('+')
                size_part = parts[0]
                original_x = int(parts[1])
                original_y = int(parts[2])
            except Exception:
                return
            def shake_animation(step=0):
                try:
                    offsets = [10, -10, 10, -10, 5, -5, 0]
                    if step < len(offsets):
                        dx = offsets[step]
                        popup.geometry(f"{size_part}+{original_x + dx}+{original_y}")
                        popup.after(50, shake_animation, step + 1)
                    elif step == len(offsets):
                        popup.after(1000, lambda: pwd_entry.configure(fg_color=original_entry_color))
                except Exception:
                    pass
            shake_animation(0)
        def unlock_master():
            pwd = pwd_var.get() or ""
            if not pwd:
                mark_wrong()
                messagebox.showerror("Error", "Master password required!", parent=popup)
                return
            try:
                key = derive_key(pwd, self.salt)
                fernet_test = Fernet(key)
                # For import, no verification yet, so skip or assume
                # For normal unlock, check verification
                if "verification" in self.settings:
                    verif_dec = fernet_test.decrypt(self.settings["verification"].encode()).decode()
                    if verif_dec != "VERIFICATION":
                        raise ValueError("Verification failed")
                self.fernet = fernet_test
                self.master_password = pwd
                closed_by_user["val"] = False
                popup.grab_release()
                popup.destroy()
            except Exception as e:
                mark_wrong()
                messagebox.showerror("Error", "Incorrect master password!", parent=popup)
        ctk.CTkButton(btn_frame, text="Unlock", command=unlock_master,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        ctk.CTkButton(btn_frame, text="Cancel", command=on_close, fg_color="#3a3a3a", width=120).pack(side="left", padx=6)
        pwd_entry.focus_set()
        center_popup(popup)
        self.wait_window(popup)
        return not closed_by_user["val"] and self.fernet is not None
    # --- Auth modal for normal unlock ---
    def _show_master_password_modal(self):
        # This is for normal unlock after setup
        return self._show_master_unlock_modal()
    def _verify_master_password(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Black Hole - Master Password")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.attributes("-topmost", True)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        verified = False
        closed_by_user = {"val": True}
        def on_close():
            popup.grab_release()
            popup.destroy()
        popup.protocol("WM_DELETE_WINDOW", on_close)
        ctk.CTkLabel(popup, text="Verify Master Password",
                     font=("Nunito", 16, "bold"),
                     text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        ctk.CTkLabel(popup, text="Enter your master password to proceed",
                     text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0,12))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        pwd_var = StringVar()
        pwd_entry = ctk.CTkEntry(frame, placeholder_text="Master Password",
                                 show="*", textvariable=pwd_var,
                                 width=360, justify="center")
        pwd_entry.pack(padx=12, pady=(12,6))
        original_entry_color = pwd_entry.cget("fg_color")
        def toggle_pwd():
            pwd_entry.configure(show="" if pwd_entry.cget("show")=="*" else "*")
        ctk.CTkButton(frame, text="Show/Hide", width=120,
                       command=toggle_pwd, fg_color=ACCENT, text_color=BG,
                       hover_color=ACCENT_DIM).pack(pady=(0,10))
        btn_frame = ctk.CTkFrame(popup, fg_color=BG, corner_radius=0)
        btn_frame.pack(pady=(12,12))
        def mark_wrong():
            pwd_entry.configure(fg_color="#7a2d2d")
            try:
                geo_string = popup.geometry()
                parts = geo_string.split('+')
                size_part = parts[0]
                original_x = int(parts[1])
                original_y = int(parts[2])
            except Exception:
                return
            def shake_animation(step=0):
                try:
                    offsets = [10, -10, 10, -10, 5, -5, 0]
                    if step < len(offsets):
                        dx = offsets[step]
                        popup.geometry(f"{size_part}+{original_x + dx}+{original_y}")
                        popup.after(50, shake_animation, step + 1)
                    elif step == len(offsets):
                        popup.after(1000, lambda: pwd_entry.configure(fg_color=original_entry_color))
                except Exception:
                    pass
            shake_animation(0)
        def verify_master():
            nonlocal verified
            pwd = pwd_var.get() or ""
            if not pwd:
                mark_wrong()
                messagebox.showerror("Error", "Master password required!", parent=popup)
                return
            try:
                key = derive_key(pwd, self.salt)
                fernet_test = Fernet(key)
                if "verification" in self.settings:
                    verif_dec = fernet_test.decrypt(self.settings["verification"].encode()).decode()
                    if verif_dec != "VERIFICATION":
                        raise ValueError("Verification failed")
                self.fernet = fernet_test
                self.master_password = pwd
                self.authenticated = True
                verified = True
                closed_by_user["val"] = False
                popup.grab_release()
                popup.destroy()
            except Exception as e:
                mark_wrong()
                messagebox.showerror("Error", "Incorrect master password!", parent=popup)
        ctk.CTkButton(btn_frame, text="Verify", command=verify_master,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=200).pack(side="left", padx=12)
        ctk.CTkButton(btn_frame, text="Cancel", command=on_close, fg_color="#3a3a3a", width=120).pack(side="left", padx=6)
        pwd_entry.focus_set()
        pwd_entry.bind("<Return>", lambda e: verify_master())
        center_popup(popup)
        self.wait_window(popup)
        return not closed_by_user["val"] and verified
    # --- Initialize DB ---
    def _init_db(self):
        self.conn = sqlite3.connect(self.db_path)
        self.c = self.conn.cursor()
        self.c.execute('''
            CREATE TABLE IF NOT EXISTS passwords (
                id INTEGER PRIMARY KEY,
                title TEXT,
                username TEXT,
                password TEXT,
                notes TEXT
            )
        ''')
        try:
            self.c.execute("ALTER TABLE passwords ADD COLUMN icon_path TEXT")
        except sqlite3.OperationalError as e:
            if "duplicate column" not in str(e).lower():
                raise
        self.conn.commit()
    # --- Save Order ---
    def _save_order(self):
        data = {
            "mode": self.order_mode,
            "custom_order": self.custom_order
        }
        try:
            with open(order_path, "w", encoding="utf-8") as f:
                json.dump(data, f)
        except Exception:
            pass
    # --- Save Pinned ---
    def _save_pinned(self):
        try:
            with open(pinned_path, "w", encoding="utf-8") as f:
                json.dump({"pinned": self.pinned}, f)
        except Exception:
            pass
    # --- Save Settings ---
    def _save_settings(self):
        try:
            with open(settings_path, "w", encoding="utf-8") as f:
                json.dump(self.settings, f)
        except Exception:
            pass
    # --- Build UI ---
    def _build_ui(self):
        header = ctk.CTkFrame(self, fg_color=BG, height=64)
        header.pack(fill="x", padx=12, pady=(12,0))
        pil_image_bh = Image.open(BLACK_HOLE_WIDE_LOGO)
        bh_width, bh_height = pil_image_bh.size
        new_width_bh = 200
        new_height_bh = int((new_width_bh / bh_width) * bh_height)
        bh_ctk_image = ctk.CTkImage(
            light_image=pil_image_bh,
            dark_image=pil_image_bh,
            size=(new_width_bh, new_height_bh)
        )
        bh_label = ctk.CTkLabel(header, image=bh_ctk_image, text="")
        bh_label.pack(side="left", padx=12, pady=(12, 6))
        # Search bar
        search_frame = ctk.CTkFrame(header, fg_color=BG)
        search_frame.pack(side="left", fill="x", expand=True)
        search_subframe = ctk.CTkFrame(search_frame, fg_color=BG)
        search_subframe.pack(pady=16, padx=20, fill="x")
        search_label = ctk.CTkLabel(search_subframe, text="🔍 Search", font=("Nunito", 12), text_color=TEXT)
        search_label.pack(side="left", padx=(0, 5))
        self.search_var = StringVar()
        self.search_entry = ctk.CTkEntry(search_subframe, textvariable=self.search_var, placeholder_text="by title")
        self.search_entry.pack(side="left", fill="x", expand=True)
        self.search_entry.bind("<KeyRelease>", lambda e: self.load_cards())
        # Sort options
        sort_frame = ctk.CTkFrame(header, fg_color=BG)
        sort_frame.pack(side="right", padx=12)
        ctk.CTkLabel(sort_frame, text="Sort:", font=("Nunito", 12), text_color=TEXT, fg_color=BG).pack(side="left", padx=4)
        sort_values = ["Default", "Title A-Z", "Title Z-A", "Custom"]
        display_mode = {
            "default": "Default",
            "a-z": "Title A-Z",
            "z-a": "Title Z-A",
            "custom": "Custom"
        }.get(self.order_mode, "Default")
        self.sort_var = StringVar(value=display_mode)
        self.sort_combo = ctk.CTkComboBox(sort_frame, values=sort_values, variable=self.sort_var, width=100, font=("Nunito", 12), command=self._change_sort)
        self.sort_combo.pack(side="left", padx=4)
        if self.order_mode == "custom":
            self.edit_order_btn = ctk.CTkButton(sort_frame, text="Edit Order", command=self.edit_custom_order,
                                                fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=80, font=("Nunito", 12))
            self.edit_order_btn.pack(side="left", padx=4)
        settings_btn = ctk.CTkButton(header, text="⚙️", command=self.show_settings_popup, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=12, font=("Nunito", 12))
        settings_btn.pack(side="right", padx=4)
        Tooltip(settings_btn, "Settings")
        add_btn = ctk.CTkButton(header, text="➕", command=self.create_new_card,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=12, font=("Nunito", 12))
        add_btn.pack(side="right", padx=4)
        Tooltip(add_btn, "Add New Entry")
        export_btn = ctk.CTkButton(header, text="📤", command=self.export_popup,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=12, font=("Nunito", 12))
        export_btn.pack(side="right", padx=4)
        Tooltip(export_btn, "Export Vault")
        import_btn = ctk.CTkButton(header, text="📥", command=self.import_spreadsheet,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=12, font=("Nunito", 12))
        import_btn.pack(side="right", padx=4)
        Tooltip(import_btn, "Import from Spreadsheet")
        about_btn = ctk.CTkButton(header, text="ℹ️", command=self.show_about,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=12, font=("Nunito", 12))
        about_btn.pack(side="right", padx=4)
        Tooltip(about_btn, "About")
        self.cards_frame = ctk.CTkScrollableFrame(self, fg_color=BG, corner_radius=10)
        self.cards_frame.pack(padx=12, pady=12, fill="both", expand=True)
        # Keyboard bindings for main window
        self.bind("<Control-f>", lambda e: self.search_entry.focus())
        self.bind("<Control-s>", lambda e: self.export_popup())
        self.bind("<Up>", lambda e: self.cards_frame._parent_canvas.yview_scroll(-20, "units"))
        self.bind("<Down>", lambda e: self.cards_frame._parent_canvas.yview_scroll(20, "units"))
        self.check_for_update()
    def import_spreadsheet(self):
        file_path = filedialog.askopenfilename(title="Select Spreadsheet", filetypes=[("Excel files", "*.xlsx"), ("CSV files", "*.csv")])
        if not file_path:
            return
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == '.xlsx':
                df = pd.read_excel(file_path)
            elif ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                messagebox.showerror("Error", "Unsupported file type")
                return
        except Exception as e:
            messagebox.showerror("Error", f"Failed to read file: {e}")
            return
        columns = list(df.columns)
        columns = [str(col) for col in columns]
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Map Columns")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.iconbitmap(APP_ICON_PATH)
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Map Spreadsheet Columns", font=("Nunito", 16, "bold"), text_color=TEXT).pack(pady=(16,6))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both")
        # Title
        ctk.CTkLabel(frame, text="Title (required):", text_color=TEXT).pack(anchor="w", padx=12, pady=(12,0))
        title_var = StringVar()
        title_combo = ctk.CTkComboBox(frame, values=columns, variable=title_var, width=360)
        title_combo.pack(padx=12, pady=(0,6))
        # Username
        ctk.CTkLabel(frame, text="Username:", text_color=TEXT).pack(anchor="w", padx=12, pady=(6,0))
        user_var = StringVar(value="None")
        user_combo = ctk.CTkComboBox(frame, values=["None"] + columns, variable=user_var, width=360)
        user_combo.pack(padx=12, pady=(0,6))
        # Password
        ctk.CTkLabel(frame, text="Password:", text_color=TEXT).pack(anchor="w", padx=12, pady=(6,0))
        pwd_var = StringVar(value="None")
        pwd_combo = ctk.CTkComboBox(frame, values=["None"] + columns, variable=pwd_var, width=360)
        pwd_combo.pack(padx=12, pady=(0,6))
        # Notes
        ctk.CTkLabel(frame, text="Notes:", text_color=TEXT).pack(anchor="w", padx=12, pady=(6,0))
        notes_var = StringVar(value="None")
        notes_combo = ctk.CTkComboBox(frame, values=["None"] + columns, variable=notes_var, width=360)
        notes_combo.pack(padx=12, pady=(0,12))
        def do_import():
            title_col = title_var.get()
            if not title_col:
                messagebox.showerror("Error", "Title column required!", parent=popup)
                return
            user_col = user_var.get() if user_var.get() != "None" else None
            pwd_col = pwd_var.get() if pwd_var.get() != "None" else None
            notes_col = notes_var.get() if notes_var.get() != "None" else None
            count = 0
            new_ids = []
            for _, row in df.iterrows():
                title = row[title_col] if title_col in row else ""
                if pd.isna(title) or not str(title).strip():
                    continue
                user = row[user_col] if user_col and user_col in row else ""
                pwd = row[pwd_col] if pwd_col and pwd_col in row else ""
                notes = row[notes_col] if notes_col and notes_col in row else ""
                enc_pwd = self.fernet.encrypt(str(pwd).encode()).decode() if pwd else ""
                self.c.execute("INSERT INTO passwords (title, username, password, notes, icon_path) VALUES (?, ?, ?, ?, ?)",
                               (str(title), str(user), enc_pwd, str(notes), ""))
                self.conn.commit()
                count += 1
                if self.order_mode == "custom":
                    new_ids.append(self.c.lastrowid)
            if self.order_mode == "custom" and new_ids:
                self.custom_order.extend(new_ids)
                self._save_order()
            popup.destroy()
            self.load_cards()
            messagebox.showinfo("Imported", f"Imported {count} entries successfully.")
        ctk.CTkButton(popup, text="Import", command=do_import, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM).pack(pady=(0,12))
        center_popup(popup)
    # --- Settings Popup ---
    def show_settings_popup(self):
        if not sys.platform.startswith("win"):
            messagebox.showinfo("Info", "Settings are only available on Windows.")
            return
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Settings")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Settings", font=("Nunito", 16, "bold"), text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        frame = ctk.CTkFrame(popup, fg_color=CARD, corner_radius=8)
        frame.pack(padx=20, pady=8, fill="both", expand=False)
        launch_var = ctk.CTkSwitch(frame, text="Launch with Windows")
        launch_var.pack(pady=10, padx=10)
        if self.settings.get("launch_with_windows", False):
            launch_var.select()
        launch_var.configure(command=lambda: self.toggle_launch(launch_var.get()))
        tray_var = ctk.CTkSwitch(frame, text="Minimize to Tray")
        tray_var.pack(pady=10, padx=10)
        if self.settings.get("minimize_to_tray", False):
            tray_var.select()
        tray_var.configure(command=lambda: self.toggle_tray(tray_var.get()))
        ctk.CTkButton(frame, text="About", command=self.show_about, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(pady=10, padx=10)
        ctk.CTkButton(frame, text="Reset", command=self.reset_app, fg_color="#ff4d4d", text_color=BG, hover_color="#ff0000", width=120).pack(pady=10, padx=10)
        ctk.CTkButton(popup, text="Close", command=lambda: popup.destroy(), fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(pady=12)
        center_popup(popup)
    def toggle_launch(self, value):
        self.settings["launch_with_windows"] = bool(value)
        self.toggle_startup(value)
        self._save_settings()
    def toggle_tray(self, value):
        self.settings["minimize_to_tray"] = bool(value)
        # Update startup if launch is enabled
        if self.settings.get("launch_with_windows", False):
            self.toggle_startup(True)
        self._save_settings()
    def toggle_startup(self, enable):
        if not sys.platform.startswith("win"):
            return
        reg_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
        app_name = "BlackHole"
        if enable:
            try:
                cmd = f'"{sys.argv[0]}"'
                if self.settings.get("minimize_to_tray", False):
                    cmd += " --minimize"
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, reg_path, 0, winreg.KEY_SET_VALUE)
                winreg.SetValueEx(key, app_name, 0, winreg.REG_SZ, cmd)
                winreg.CloseKey(key)
            except Exception as e:
                messagebox.showerror("Error", f"Failed to add to startup: {e}")
        else:
            try:
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, reg_path, 0, winreg.KEY_SET_VALUE)
                winreg.DeleteValue(key, app_name)
                winreg.CloseKey(key)
            except FileNotFoundError:
                pass
            except Exception as e:
                messagebox.showerror("Error", f"Failed to remove from startup: {e}")
    def _change_sort(self, event=None):
        val = self.sort_var.get()
        if val == "Title A-Z":
            self.order_mode = "a-z"
        elif val == "Title Z-A":
            self.order_mode = "z-a"
        elif val == "Custom":
            self.order_mode = "custom"
        else:
            self.order_mode = "default"
        self._save_order()
        self.load_cards()
        if self.order_mode == "custom":
            if not hasattr(self, "edit_order_btn"):
                sort_frame = self.sort_combo.master
                self.edit_order_btn = ctk.CTkButton(sort_frame, text="Edit Order", command=self.edit_custom_order,
                                                    fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=80, font=("Nunito", 12))
                self.edit_order_btn.pack(side="left", padx=4)
        else:
            if hasattr(self, "edit_order_btn"):
                self.edit_order_btn.destroy()
                del self.edit_order_btn
    # --- Load Cards ---
    def load_cards(self):
        for widget in self.cards_frame.winfo_children():
            widget.destroy()
        search_term = self.search_var.get().strip().lower()
        if self.order_mode == "default":
            self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords ORDER BY id DESC")
        elif self.order_mode == "a-z":
            self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords ORDER BY LOWER(title) ASC")
        elif self.order_mode == "z-a":
            self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords ORDER BY LOWER(title) DESC")
        else: # custom
            self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords")
            rows = self.c.fetchall()
            current_ids = [r[0] for r in rows]
            self.custom_order = [id_ for id_ in self.custom_order if id_ in current_ids]
            new_ids = [id_ for id_ in current_ids if id_ not in self.custom_order]
            new_ids.sort(reverse=True)
            self.custom_order += new_ids
            self._save_order()
            id_to_row = {r[0]: r for r in rows}
            rows = [id_to_row[id_] for id_ in self.custom_order if id_ in id_to_row]
        rows = self.c.fetchall() if self.order_mode != "custom" else rows
        if search_term:
            rows = [r for r in rows if search_term in r[1].lower()]
        match(self.winfo_screenwidth()):
            case _ if self.winfo_screenwidth() >= 2200:
                num_columns = 6
            case _ if self.winfo_screenwidth() >= 1900:
                num_columns = 5
            case _ if self.winfo_screenwidth() >= 1500:
                num_columns = 4
            case _ if self.winfo_screenwidth() >= 1200:
                num_columns = 3
            case _ if self.winfo_screenwidth() >= 800:
                num_columns = 2
            case _:
                num_columns = 1
        self.cards_frame.grid_columnconfigure(tuple(range(num_columns)), weight=1)
        for i, row in enumerate(rows):
            id_, title, user, pwd_enc, notes, _ = row # Ignore stored icon_path
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            shadow_frame = ctk.CTkFrame(self.cards_frame, fg_color="gray20", width=360, height=470, border_width=0, corner_radius=0)
            row_num = i // num_columns
            col = i % num_columns
            shadow_frame.grid(row=row_num, column=col, padx=10, pady=10, sticky="n")
            card = ctk.CTkFrame(shadow_frame, fg_color=CARD, corner_radius=0, width=360, height=450, border_width=0)
            card.grid(padx=8, pady=8, sticky="nsew") # Use grid with padding for shadow effect
            shadow_frame.grid_columnconfigure(0, weight=1)
            shadow_frame.grid_rowconfigure(0, weight=1)
            image_frame = ctk.CTkFrame(card, height=350, fg_color="transparent", corner_radius=0)
            image_frame.pack(fill="x", expand=False)
            # Search for icon by file name
            icon_path = None
            for ext in ['.png', '.jpg', '.jpeg']:
                possible_path = os.path.join(stored_icons_path, f"{id_}{ext}")
                if os.path.exists(possible_path):
                    icon_path = possible_path
                    break
            if icon_path:
                try:
                    pil_img = Image.open(icon_path)
                    ctk_img = ctk.CTkImage(light_image=pil_img, size=(350, 350))
                    image_label = ctk.CTkLabel(image_frame, image=ctk_img, text="")
                    image_label.pack(expand=True, fill="both")
                except:
                    pass
            bottom_frame = ctk.CTkFrame(card, height=150, fg_color=CARD)
            bottom_frame.pack(fill="x", expand=True)
            left = ctk.CTkFrame(bottom_frame, fg_color=CARD, corner_radius=0)
            left.pack(side="left", fill="both", expand=True, padx=4, pady=4)
            title_label = ctk.CTkLabel(left, text=title or "(No title)", anchor="w",
                        font=("Nunito", 14, "bold"), text_color=TEXT, fg_color=CARD, width=50, wraplength=200) # Increased font size
            title_label.pack(anchor="w")
            user_label = ctk.CTkLabel(left, text=f"User: {user or ''}", anchor="w",
                        text_color=ACCENT_DIM, fg_color=CARD, font=("Nunito", 12), width=50, wraplength=200) # Increased font size
            user_label.pack(anchor="w")
            pwd_var = StringVar(value="*"*len(pwd) if pwd else "")
            pwd_label = ctk.CTkLabel(left, textvariable=pwd_var, anchor="w", text_color=TEXT, fg_color=CARD, font=("Nunito", 12), width=50, wraplength=200) # Increased font size
            pwd_label.pack(anchor="w")
            def copy_text(text, msg="Copied to clipboard!"):
                if text:
                    self.clipboard_clear()
                    self.clipboard_append(text)
                    messagebox.showinfo("Copied", msg)
            user_label.bind("<Button-1>", lambda e, u=user: copy_text(u, "Username copied!"))
            pwd_label.bind("<Button-1>", lambda e, p=pwd: copy_text(p, "Password copied!"))
            right = ctk.CTkFrame(bottom_frame, fg_color=CARD, corner_radius=0)
            right.pack(side="right", padx=4, pady=4)
            righter = ctk.CTkFrame(right, fg_color=CARD, corner_radius=0)
            righter.pack(side="right", padx=4, pady=0)
            def toggle_show(pw=pwd, var=pwd_var):
                var.set(pw if var.get().startswith("*") else "*"*len(pw))
            ctk.CTkButton(right, text="Show", command=toggle_show,
                        width=60, fg_color=ACCENT, text_color=BG, font=("Nunito", 10)).pack(pady=2) # Increased font size
            ctk.CTkButton(right, text="Edit", command=lambda id=id_: self.edit_card_popup(id), width=60, font=("Nunito", 10)).pack(pady=2) # Increased font size
            def show_notes(n=notes):
                messagebox.showinfo("Notes", n or "No notes")
            ctk.CTkButton(righter, text="Notes", command=lambda n=notes: show_notes(n), width=60, font=("Nunito", 10)).pack(pady=2) # Added Notes button
            ctk.CTkButton(righter, text="Delete", command=lambda id=id_: self.delete_card(id), width=60,
                        fg_color="#7a2d2d", font=("Nunito", 10)).pack(pady=2) # Increased font size
    # --- Create Card ---
    def create_new_card(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Create New Password")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Create New Entry", font=("Nunito", 14, "bold"), text_color=TEXT, fg_color=BG, justify="center").pack(pady=(12,6))
        title_var = StringVar()
        title_entry = ctk.CTkEntry(popup, placeholder_text="Title (required)", textvariable=title_var, width=320)
        title_entry.pack(pady=8)
        def create_card_action():
            title = title_var.get().strip()
            if not title:
                messagebox.showerror("Error", "Title required!", parent=popup)
                return
            self.c.execute("INSERT INTO passwords (title, username, password, notes, icon_path) VALUES (?, ?, ?, ?, ?)", (title, "", "", "", ""))
            self.conn.commit()
            if self.order_mode == "custom":
                new_id = self.c.lastrowid
                self.custom_order.append(new_id)
                self._save_order()
            popup.grab_release()
            popup.destroy()
            self.load_cards()
        ctk.CTkButton(popup, text="Create", command=create_card_action,
                       fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=140).pack(pady=(12,10))
        title_entry.focus_set()
        popup.bind("<Return>", lambda e: create_card_action())
        center_popup(popup)
    # --- Edit Card Popup ---
    def edit_card_popup(self, id_):
        row = self.c.execute("SELECT title, username, password, notes, icon_path FROM passwords WHERE id=?", (id_,)).fetchone()
        if not row:
            messagebox.showerror("Error", "Entry not found.")
            return
        title, user, pwd_enc, notes, current_icon_path = row
        try:
            pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
        except Exception:
            pwd = ""
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Edit Entry")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Edit Entry", font=("Nunito", 14, "bold"), text_color=TEXT, fg_color=BG).pack(pady=(12,6))
        ctk.CTkLabel(popup, text="Title", text_color=ACCENT_DIM, fg_color=BG).pack(anchor="w", padx=20)
        title_var = StringVar(value=title)
        ctk.CTkEntry(popup, textvariable=title_var, width=420).pack(padx=20, pady=(4,8))
        ctk.CTkLabel(popup, text="Username", text_color=ACCENT_DIM, fg_color=BG).pack(anchor="w", padx=20)
        user_var = StringVar(value=user)
        ctk.CTkEntry(popup, textvariable=user_var, width=420).pack(padx=20, pady=(4,8))
        ctk.CTkLabel(popup, text="Password", text_color=ACCENT_DIM, fg_color=BG).pack(anchor="w", padx=20)
        pwd_var = StringVar(value=pwd)
        pwd_entry = ctk.CTkEntry(popup, textvariable=pwd_var, show="*", width=420)
        pwd_entry.pack(padx=20, pady=(4,8))
        def toggle_pwd_entry():
            pwd_entry.configure(show="" if pwd_entry.cget("show")=="*" else "*")
        ctk.CTkButton(popup, text="Show/Hide", command=toggle_pwd_entry, fg_color=ACCENT, text_color=BG).pack(pady=(0,8))
        ctk.CTkLabel(popup, text="Notes", text_color=ACCENT_DIM, fg_color=BG).pack(anchor="w", padx=20)
        notes_var = StringVar(value=notes)
        ctk.CTkEntry(popup, textvariable=notes_var, width=420).pack(padx=20, pady=(4,12))
        ctk.CTkLabel(popup, text="Icon", text_color=ACCENT_DIM, fg_color=BG).pack(anchor="w", padx=20)
        new_icon_path = None
        def upload_icon():
            nonlocal new_icon_path
            file = filedialog.askopenfilename(filetypes=[("Images", "*.png *.jpg *.jpeg *.ico")])
            if file:
                try:
                    pil_img = Image.open(file)
                    pil_img = pil_img.resize((350, 350), Image.LANCZOS)
                    orig_ext = os.path.splitext(file)[1].lower()
                    if orig_ext == '.ico':
                        orig_ext = '.png'
                    dest = os.path.join(stored_icons_path, f"{id_}{orig_ext}")
                    format_to_save = 'JPEG' if orig_ext == '.jpg' or orig_ext == '.jpeg' else 'PNG'
                    pil_img.save(dest, format=format_to_save)
                    # Remove old icons with different extensions
                    for other_ext in ['.png', '.jpg', '.jpeg']:
                        if other_ext != orig_ext:
                            old_path = os.path.join(stored_icons_path, f"{id_}{other_ext}")
                            if os.path.exists(old_path):
                                os.remove(old_path)
                    new_icon_path = dest
                    messagebox.showinfo("Uploaded", "Icon uploaded and resized!", parent=popup)
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to process icon: {e}", parent=popup)
        ctk.CTkButton(popup, text="Upload Icon", command=upload_icon, fg_color=ACCENT, text_color=BG).pack(pady=(0,8))
        pinned_var = ctk.CTkSwitch(popup, text="Pinned to Tray")
        pinned_var.pack(pady=8)
        if id_ in self.pinned:
            pinned_var.select()
        def save_card():
            enc_pwd = self.fernet.encrypt(pwd_var.get().encode()).decode() if pwd_var.get() else ""
            icon_to_save = new_icon_path if new_icon_path else current_icon_path
            if new_icon_path and current_icon_path and os.path.exists(current_icon_path):
                try:
                    os.remove(current_icon_path)
                except:
                    pass
            self.c.execute("UPDATE passwords SET title=?, username=?, password=?, notes=?, icon_path=? WHERE id=?",
                           (title_var.get(), user_var.get(), enc_pwd, notes_var.get(), icon_to_save or "", id_))
            self.conn.commit()
            if pinned_var.get():
                if id_ not in self.pinned:
                    self.pinned.append(id_)
            else:
                if id_ in self.pinned:
                    self.pinned.remove(id_)
            self._save_pinned()
            popup.grab_release()
            popup.destroy()
            self.load_cards()
        ctk.CTkButton(popup, text="Save", command=save_card, fg_color=ACCENT, text_color=BG, width=120).pack(pady=(6,12))
        popup.bind("<Control-s>", lambda e: save_card())
        popup.bind("<Return>", lambda e: save_card())
        center_popup(popup)
    # --- Delete ---
    def delete_card(self, id_):
        if messagebox.askyesno("Confirm Delete", "Are you sure you want to delete this entry?"):
            row = self.c.execute("SELECT icon_path FROM passwords WHERE id=?", (id_,)).fetchone()
            if row and row[0] and os.path.exists(row[0]):
                try:
                    os.remove(row[0])
                except:
                    pass
            # Also remove any searched icon files
            for ext in ['.png', '.jpg', '.jpeg']:
                possible_path = os.path.join(stored_icons_path, f"{id_}{ext}")
                if os.path.exists(possible_path):
                    try:
                        os.remove(possible_path)
                    except:
                        pass
            self.c.execute("DELETE FROM passwords WHERE id=?", (id_,))
            self.conn.commit()
            if self.order_mode == "custom" and id_ in self.custom_order:
                self.custom_order.remove(id_)
                self._save_order()
            if id_ in self.pinned:
                self.pinned.remove(id_)
                self._save_pinned()
            self.load_cards()
    # --- Custom Order Popup ---
    def edit_custom_order(self):
        # Update custom_order with current entries
        self.c.execute("SELECT id FROM passwords")
        current_ids = [r[0] for r in self.c.fetchall()]
        self.custom_order = [id_ for id_ in self.custom_order if id_ in current_ids]
        new_ids = [id_ for id_ in current_ids if id_ not in self.custom_order]
        new_ids.sort(reverse=True)
        self.custom_order += new_ids
        self._save_order()
        # Get titles
        self.c.execute("SELECT id, title FROM passwords")
        id_to_title = {r[0]: r[1] for r in self.c.fetchall()}
        popup = ctk.CTkToplevel(self)
        popup.title("Custom Order")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        popup.grab_set()
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Reorder Entries", font=("Nunito", 16, "bold"), text_color=TEXT, fg_color=BG).pack(pady=(16,6))
        frame = ctk.CTkFrame(popup, fg_color=BG)
        frame.pack(padx=20, pady=10, fill="both", expand=True)
        lb = Listbox(frame, bg=CARD, fg=TEXT, selectbackground=ACCENT, font=("Nunito", 12), height=15, width=40)
        lb.pack(side="left", fill="both", expand=True)
        scroll = ctk.CTkScrollbar(frame, command=lb.yview)
        scroll.pack(side="right", fill="y")
        lb.configure(yscrollcommand=scroll.set)
        reorder_ids = self.custom_order[:]
        for id_ in reorder_ids:
            title = id_to_title.get(id_, f"ID {id_}")
            lb.insert(END, title or "(No title)")
        def move_up(event=None):
            try:
                i = lb.curselection()[0]
                if i == 0:
                    return "break"
                text = lb.get(i)
                lb.delete(i)
                lb.insert(i-1, text)
                lb.selection_set(i-1)
                # Swap ids
                reorder_ids[i], reorder_ids[i-1] = reorder_ids[i-1], reorder_ids[i]
                return "break"
            except:
                pass
            return "break"
        def move_down(event=None):
            try:
                i = lb.curselection()[0]
                if i == lb.size() - 1:
                    return "break"
                text = lb.get(i)
                lb.delete(i)
                lb.insert(i+1, text)
                lb.selection_set(i+1)
                # Swap ids
                reorder_ids[i], reorder_ids[i+1] = reorder_ids[i+1], reorder_ids[i]
                return "break"
            except:
                pass
            return "break"
        lb.bind("<Up>", move_up)
        lb.bind("<Down>", move_down)
        btn_frame = ctk.CTkFrame(popup, fg_color=BG)
        btn_frame.pack(pady=10)
        ctk.CTkButton(btn_frame, text="Up", command=move_up, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=60).pack(side="left", padx=5)
        ctk.CTkButton(btn_frame, text="Down", command=move_down, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=60).pack(side="left", padx=5)
        def save_reorder():
            self.custom_order = reorder_ids
            self._save_order()
            self.load_cards()
            popup.grab_release()
            popup.destroy()
        ctk.CTkButton(popup, text="Save", command=save_reorder, fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(pady=(0,12))
        popup.bind("<Control-s>", lambda e: save_reorder())
        popup.bind("<Return>", lambda e: save_reorder())
        center_popup(popup)
    # --- Export ---
    def export_popup(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("Export Passwords")
        popup.configure(fg_color=BG)
        popup.geometry("1000x500")
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(popup, text="Export Passwords", font=("Nunito", 14, "bold"), text_color=TEXT, fg_color=BG).pack(pady=(12,6))
        export_frame = ctk.CTkScrollableFrame(popup, fg_color=BG, width=400, height=210)
        export_frame.pack(pady=12, padx=12, fill="both", expand=True)
        def add_section(parent, name):
            section = ctk.CTkFrame(parent, fg_color=CARD, corner_radius=8)
            section.pack(fill="x", pady=10, padx=10)
            ctk.CTkLabel(section, text=name, font=("Nunito", 16, "bold"), text_color=TEXT).pack(pady=5)
            btn_frame = ctk.CTkFrame(section, fg_color=CARD)
            btn_frame.pack(pady=5)
            return btn_frame
        doc_btns = add_section(export_frame, "Documents")
        ctk.CTkButton(doc_btns, text=".docx", command=self.export_docx,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(doc_btns, text=".odt", command=self.export_odt,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(doc_btns, text=".txt", command=self.export_txt,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        sheet_btns = add_section(export_frame, "Spreadsheets")
        ctk.CTkButton(sheet_btns, text=".xlsx", command=self.export_xlsx,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(sheet_btns, text=".ods", command=self.export_ods,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(sheet_btns, text=".csv", command=self.export_csv,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        slide_btns = add_section(export_frame, "Slides")
        ctk.CTkButton(slide_btns, text=".odp", command=self.export_odp,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(slide_btns, text=".pptx", command=self.export_pptx,
                      fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(side="left", padx=5)
        ctk.CTkButton(popup, text="Cancel", command=popup.destroy, fg_color="#3a3a3a", width=120).pack(pady=12)
        center_popup(popup)
        self.wait_window(popup)
    def export_docx(self):
        if not self._verify_master_password():
            return
        doc = Document()
        for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
            title, user, pwd_enc, notes = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            doc.add_paragraph(f"Title: {title}")
            doc.add_paragraph(f"Username: {user}")
            doc.add_paragraph(f"Password: {pwd}")
            doc.add_paragraph(f"Notes: {notes}")
            doc.add_paragraph("-"*30)
        path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word","*.docx")])
        if path:
            doc.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_odt(self):
        if not self._verify_master_password():
            return
        odt = OpenDocumentText()
        for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
            title, user, pwd_enc, notes = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            for text in [f"Title: {title}", f"Username: {user}", f"Password: {pwd}", f"Notes: {notes}", "-"*30]:
                p = P(text=text)
                odt.text.addElement(p)
        path = filedialog.asksaveasfilename(defaultextension=".odt", filetypes=[("OpenDocument","*.odt")])
        if path:
            odt.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_txt(self):
        if not self._verify_master_password():
            return
        path = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text","*.txt")])
        if path:
            with open(path, "w", encoding="utf-8") as f:
                for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
                    title, user, pwd_enc, notes = row
                    try:
                        pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
                    except Exception:
                        pwd = ""
                    f.write(f"Title: {title}\n")
                    f.write(f"Username: {user}\n")
                    f.write(f"Password: {pwd}\n")
                    f.write(f"Notes: {notes}\n")
                    f.write("-"*30 + "\n")
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_xlsx(self):
        if not self._verify_master_password():
            return
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Title", "Username", "Password", "Notes"])
        for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
            title, user, pwd_enc, notes = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            ws.append([title, user, pwd, notes])
        path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel","*.xlsx")])
        if path:
            wb.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_ods(self):
        if not self._verify_master_password():
            return
        doc = OpenDocumentSpreadsheet()
        table = Table(name="Passwords")
        doc.spreadsheet.addElement(table)
        header_row = TableRow()
        for header in ["Title", "Username", "Password", "Notes"]:
            cell = TableCell()
            p = P(text=header)
            cell.addElement(p)
            header_row.addElement(cell)
        table.addElement(header_row)
        for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
            title, user, pwd_enc, notes = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            data_row = TableRow()
            for value in [title, user, pwd, notes]:
                cell = TableCell()
                p = P(text=value)
                cell.addElement(p)
                data_row.addElement(cell)
            table.addElement(data_row)
        path = filedialog.asksaveasfilename(defaultextension=".ods", filetypes=[("OpenDocument Spreadsheet","*.ods")])
        if path:
            doc.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_csv(self):
        if not self._verify_master_password():
            return
        data = []
        for row in self.c.execute("SELECT title, username, password, notes FROM passwords"):
            title, user, pwd_enc, notes = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            data.append({"Title": title, "Username": user, "Password": pwd, "Notes": notes})
        df = pd.DataFrame(data)
        path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV","*.csv")])
        if path:
            df.to_csv(path, index=False)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_odp(self):
        if not self._verify_master_password():
            return
        doc = OpenDocumentPresentation()
        master = MasterPage(name="Default")
        doc.masterstyles.addElement(master)
        rows = self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords").fetchall()
        for row in rows:
            id_, title, user, pwd_enc, notes, icon_path = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            page = Page(masterpagename="Default")
            doc.presentation.addElement(page)
            if icon_path and os.path.exists(icon_path):
                href = doc.addPicture(icon_path)
                image_frame = Frame(width="10cm", height="10cm", x="2cm", y="4cm", anchortype="page")
                image = Image(href=href, type="simple", show="embed", actuate="onLoad")
                image_frame.addElement(image)
                page.addElement(image_frame)
            text_frame = Frame(width="15cm", height="10cm", x="14cm", y="4cm", anchortype="page")
            textbox = TextBox()
            text_frame.addElement(textbox)
            textbox.addElement(P(text=f"Title: {title}"))
            textbox.addElement(P(text=f"Username: {user}"))
            textbox.addElement(P(text=f"Password: {pwd}"))
            textbox.addElement(P(text=f"Notes: {notes}"))
            page.addElement(text_frame)
        path = filedialog.asksaveasfilename(defaultextension=".odp", filetypes=[("OpenDocument Presentation","*.odp")])
        if path:
            doc.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    def export_pptx(self):
        if not self._verify_master_password():
            return
        prs = Presentation()
        rows = self.c.execute("SELECT id, title, username, password, notes, icon_path FROM passwords").fetchall()
        for row in rows:
            id_, title, user, pwd_enc, notes, icon_path = row
            try:
                pwd = self.fernet.decrypt(pwd_enc.encode()).decode() if pwd_enc else ""
            except Exception:
                pwd = ""
            slide_layout = prs.slide_layouts[6] # blank slide
            slide = prs.slides.add_slide(slide_layout)
            if icon_path and os.path.exists(icon_path):
                slide.shapes.add_picture(icon_path, Inches(1), Inches(1), Inches(5), Inches(5))
            txBox = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(5), Inches(5))
            tf = txBox.text_frame
            tf.add_paragraph().text = f"Title: {title}"
            tf.add_paragraph().text = f"Username: {user}"
            tf.add_paragraph().text = f"Password: {pwd}"
            tf.add_paragraph().text = f"Notes: {notes}"
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint","*.pptx")])
        if path:
            prs.save(path)
            messagebox.showinfo("Exported", f"Exported to {path}")
    #--- About Popup ---
    def show_about(self):
        popup = ctk.CTkToplevel(self)
        popup.grab_set()
        popup.title("About Black Hole")
        popup.configure(fg_color=BG)
        popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                popup.after(250, lambda: popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        pil_image_bh = Image.open(BLACK_HOLE_LOGO)
        bh_width, bh_height = pil_image_bh.size
        new_width_bh = 200
        new_height_bh = int((new_width_bh / bh_width) * bh_height)
        bh_ctk_image = ctk.CTkImage(
            light_image=pil_image_bh,
            dark_image=pil_image_bh,
            size=(new_width_bh, new_height_bh)
        )
        bh_label = ctk.CTkLabel(popup, image=bh_ctk_image, text="")
        bh_label.pack(pady=(12, 6))
        pil_image_nf = Image.open(NOVA_FOUNDRY_LOGO)
        nf_width, nf_height = pil_image_nf.size
        new_width_nf = 100
        new_height_nf = int((new_width_nf / nf_width) * nf_height)
        nf_ctk_image = ctk.CTkImage(
            light_image=pil_image_nf,
            dark_image=pil_image_nf,
            size=(new_width_nf, new_height_nf)
        )
        nf_label = ctk.CTkLabel(popup, image=nf_ctk_image, text="")
        nf_label.pack(pady=(0, 12))
        ctk.CTkLabel(popup, text="Black Hole Password Manager", font=("Nunito", 16, "bold"),
                    text_color=TEXT, fg_color=BG).pack(pady=(12, 6))
        ctk.CTkLabel(popup, text=f"Version {VERSION}\n\n",
                    text_color=ACCENT_DIM, fg_color=BG).pack(pady=(0, 12))
        try:
            with open(LICENSE_TEXT, 'r', encoding='utf-8') as f:
                license_content = f.read()
        except Exception as e:
            print(f"Error loading license file: {e}")
            license_content = "Could not load license information."
        license_box = ctk.CTkTextbox(popup,
                                    width=480,
                                    height=250,
                                    text_color=TEXT,
                                    fg_color=BG,
                                    wrap="word")
        license_box.insert("1.0", license_content)
        license_box.configure(state="disabled")
        license_box.pack(padx=20, pady=(0, 12))
        support_link = ctk.CTkLabel(popup, text="Support Nova Foundry", font=("Nunito", 12, "underline"),
                                    text_color=ACCENT, fg_color=BG, cursor="hand2")
        support_link.pack(pady=(0, 12))
        def open_support_link(event):
            webbrowser.open_new("https://buymeacoffee.com/novafoundry")
        support_link.bind("<Button-1>", open_support_link)
        ctk.CTkButton(popup, text="OK", command=lambda: (popup.grab_release(), popup.destroy()),
                    fg_color=ACCENT, text_color=BG, hover_color=ACCENT_DIM, width=120).pack(pady=(0, 12))
        center_popup(popup)
        self.wait_window(popup)
    def reset_app(self):
        if not messagebox.askyesno("Confirm Reset", "Are you sure you want to reset the app? This will delete settings and stored icons, but not the database."):
            return
        if not self._verify_master_password():
            return
        self.conn.close()
        shutil.rmtree(nova_folder)
        messagebox.showinfo("Reset Complete", "App reset. Restart the application.")
        self.destroy()
        sys.exit()
    def check_for_update(self):
        q = queue.Queue()
        def check_task():
            try:
                url = "https://api.github.com/repos/DirectedHunt42/BlackHole/releases/latest"
                req = urllib.request.Request(url, headers={'User-Agent': 'EchoHub', 'Accept': 'application/vnd.github.v3+json'})
                with urllib.request.urlopen(req) as response:
                    data = json.loads(response.read().decode('utf-8'))
                q.put(data)
            except:
                q.put(None)
        threading.Thread(target=check_task, daemon=True).start()
        def process_queue():
            try:
                data = q.get_nowait()
                if data:
                    self.do_update_confirm(data)
            except queue.Empty:
                pass
            self.after(100, process_queue)
        self.after(100, process_queue)
    def do_update_confirm(self, data):
        try:
            title = data.get('name', '').strip()
            if title.lower().startswith("release "):
                new_ver = title[len("Release "):].strip()
            elif title.lower().startswith("v"):
                new_ver = title[1:].strip()
            else:
                new_ver = title
            current_ver = VERSION
            def version_to_tuple(v):
                return tuple(map(int, v.strip("v").split(".")))
            if version_to_tuple(new_ver) > version_to_tuple(current_ver):
                if messagebox.askyesno("Update Available", f"A new version ({new_ver}) is available. Do you want to download and install it?"):
                    self.download_and_install(data)
        except Exception as e:
            print(f"Update check failed: {e}")
    def download_and_install(self, data):
        progress_popup = ctk.CTkToplevel(self)
        progress_popup.grab_set()
        progress_popup.title("Downloading Update")
        progress_popup.configure(fg_color=BG)
        progress_popup.resizable(False, False)
        if os.path.exists(APP_ICON_PATH):
            try:
                progress_popup.after(250, lambda: progress_popup.iconbitmap(APP_ICON_PATH))
            except Exception:
                pass
        ctk.CTkLabel(progress_popup, text="Downloading update...", font=("Nunito", 14, "bold"), text_color=TEXT, fg_color=BG).pack(pady=(12,12))
        center_popup(progress_popup)
        download_bar = ctk.CTkProgressBar(progress_popup, mode="indeterminate", width=300)
        download_bar.pack(pady=(0,12))
        download_bar.start()
        progress_popup.update()
        try:
            assets = data.get('assets', [])
            download_url = None
            for asset in assets:
                if asset.get('name','') == 'Black_hole_setup.exe':
                    download_url = asset.get('browser_download_url', None)
                    break
            if not download_url:
                raise Exception("No matching asset found")
            temp_path = os.path.join(os.getenv("TEMP") or ".", "Black_hole_setup.exe")
            req = urllib.request.Request(download_url, headers={'User-Agent': 'EchoHub'})
            with urllib.request.urlopen(req) as response, open(temp_path, 'wb') as out_file:
                shutil.copyfileobj(response, out_file)
            download_bar.stop()
            progress_popup.destroy()
            os.startfile(temp_path)
            self.quit()
        except Exception as e:
            download_bar.stop()
            progress_popup.destroy()
            messagebox.showerror("Error", f"Failed to download update: {str(e)}")
# --- Run App ---
if __name__ == "__main__":
    app = PasswordManager()
    print("Black Hole Password Manager started.")
    app.mainloop()